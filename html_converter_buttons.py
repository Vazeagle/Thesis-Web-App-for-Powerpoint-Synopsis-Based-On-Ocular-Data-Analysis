#from _typeshed import FileDescriptorLike
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor

import aspose.slides as asp_slides
import bs4

import os, shutil
#import glob # for filename opening

def html_slides_creator(powerpointFile, fileName):
    cwd = os.getcwd()
    print("Sheeesh CWD: ",cwd)
    prsDir = cwd+'/tempPowerpointSlides/'

    prs = Presentation(powerpointFile)
    #slide index starts from 0
    slide_counter = 0
    slides_sum = len(prs.slides)
    tempNames=[]
    ##xml_slides.remove(slides[3])

    while slide_counter < slides_sum:
        prs = Presentation(powerpointFile)
        xml_slides = prs.slides._sldIdLst
        slides = list((xml_slides))
        
        for i in range(slides_sum-1,slide_counter,-1):
            xml_slides.remove(slides[i])

        for i in range(0,slide_counter,1):
            xml_slides.remove(slides[i])
        prs.save(prsDir+fileName+"_"+ str(slide_counter) +'.pptx')
        slide_counter+=1
        tempNames.append(fileName+"_"+ str(slide_counter-1) +'.pptx')
        print("tempNames: ",tempNames)


    # Instantiate a Presentation object that represents a PPT file
    #Presentations slides dir
    pptxFileCounter = len(tempNames)-1

    try:
        for file in tempNames:  #tempNames has the .pptx extension with it
            print("FILE: ", file)
            if (os.path.isfile(prsDir+file)):
                presentation = asp_slides.Presentation(prsDir+file) 

                filename=file.rsplit('.', 1)[0]#remove the .pptx extension of the file for the correct filename
                print("CORECT FILENAME: ", filename)

                saveHTML=cwd+"/templates/slides/"+fileName+"/"+filename+".html"
                path=cwd+"/templates/slides/"+fileName+"/"

                #if file not exists create it!
                if not os.path.exists(path):
                    os.mkdir(path)
                    print("Directory " , path ,  " Created ")
                else:    
                    print("Directory " , path ,  " already exists")

                presentation.save(saveHTML,  asp_slides.export.SaveFormat.HTML)

                # load the file
                with open(saveHTML,"r",encoding="cp437", errors='ignore') as htmlFile: #or utf-8 with html5lib and no delete empty par but delete text of body
                    txt = htmlFile.read()
                    soup = bs4.BeautifulSoup(txt, "lxml")
                    watermark = soup.find('text',{'x':'388.125','y':'229.78125'})   #remove watermark by making it invisible IF we want we can actually delete it
                    watermark['fill-opacity'] = '0.0'
                    scale = soup.find('svg')

                    #This resolution will scale for every user that tries load the slide based on the screen resolution that the browser is currently open
                    scale['width']="1920"   #standard resolution for most pc's monitors
                    scale['height']="1080"
                    scale['id']="svg_scale"

                    script_tag = soup.new_tag("script")
                    script_tag["src"]="{{url_for('static', filename='js/slideRescale.js')}}"

                    par = soup.find('p')
                    par.decompose() #delete empty par base on lxml encoding!
                    trash = soup.find('&#xFEFF; ')
                    body = soup.body #delete margin to have correct resolution!
                    body.append(soup.new_tag('style', type='text/css'))
                    body.style.append('body {margin:0;padding:0}')
                    body["onload"]="rescale();" #rescale js function according to the resolution of the user!
                    body.insert(1,script_tag)#add script location before body of html! 



                    svg_canvas_group = soup.find('g',{'pointer-events':'painted'})

                    svg_canvas_group.append(soup.new_tag('foreignObject'))
                    foreignObject = soup.foreignObject
                    foreignObject["x"]="800"
                    foreignObject["y"]="520"
                    foreignObject["width"]="150"
                    foreignObject["height"]="25"

                    #Go to next slide
                    #get current slide number
                    print("fileName: ",fileName+"_")
                    curSlide = int(filename.replace(fileName+"_",""))
                    print("curSlide:",curSlide)
                    if curSlide == pptxFileCounter: #if its the last slide return to start
                        prevSlide = curSlide-1
                        nextLoc = fileName+"_" + "0" + ".html" 
                        prevLoc = fileName+"_" + str(prevSlide) + ".html"
                    elif curSlide==0:#if its the first slide no back stay the same
                        nextSlide = curSlide+1
                        nextLoc = fileName+"_" + str(nextSlide) + ".html"
                        prevLoc = fileName+"_" + "0" + ".html"
                    else:
                        nextSlide = curSlide+1
                        prevSlide = curSlide-1
                        nextLoc = fileName+"_" + str(nextSlide) + ".html"
                        prevLoc = fileName+"_" + str(prevSlide) + ".html"

                    btnExitSlide=soup.new_tag('a',role='button')
                    btnExitSlide["href"]="/synopsis/canceled/cancel"
                    btnExitSlide.append(soup.new_tag('style', type='text/css'))
                    exit_img = soup.new_tag('img', src="{{ url_for('static', filename='exit.png') }}", width="20", height="20")
                    btnExitSlide.append(exit_img)
                    foreignObject.append(btnExitSlide)

                    if curSlide != 0: #IF ITS NOT THE FIRST SLIDE
                        btnPrevSlide=soup.new_tag('a', role='button')
                        btnPrevSlide["href"]="/slides/"+fileName+"/"+ prevLoc+"/bck"
                        btnPrevSlide.append(soup.new_tag('style', type='text/css'))
                        btnPrevSlide.style.append('a {margin:2;}')# if i leave it as button it will interact with all button named elements if i dont want htat i should change the tag name above at 109 and here
                        prev_img = soup.new_tag('img', src="{{ url_for('static', filename='prev.png') }}", width="20", height="20")
                        btnPrevSlide.append(prev_img)
                        foreignObject.append(btnPrevSlide)

                    if curSlide != pptxFileCounter: #IF ITS NOT THE LAST SLIDE
                        btnNextSlide=soup.new_tag('a',role='button')
                        btnNextSlide["href"]="/slides/"+fileName+"/"+ nextLoc+"/fwd"
                        btnNextSlide.append(soup.new_tag('style', type='text/css'))
                        next_img = soup.new_tag('img', src="{{ url_for('static', filename='next.png') }}", width="20", height="20")
                        btnNextSlide.append(next_img)
                        foreignObject.append(btnNextSlide)

                    if curSlide == pptxFileCounter: #if its the last slide add the start synopsis button!
                        btnDoneSlide=soup.new_tag('a',role='button')
                        btnDoneSlide["href"]="/synopsis/"+fileName+"_"+str(curSlide)+"/start"
                        btnDoneSlide.append(soup.new_tag('style', type='text/css'))
                        check_img = soup.new_tag('img', src="{{ url_for('static', filename='check.png') }}", width="20", height="20")
                        btnDoneSlide.append(check_img)
                        foreignObject.append(btnDoneSlide)

                    sldTitle = soup.findAll('div',{'class':'slideTitle'})
                    for title in sldTitle:
                        title.decompose() #remove slide the extra title div


                with open(saveHTML,"w",encoding="cp437", errors='ignore') as htmlFile:
                    htmlFile.write(str(soup))
    except Exception as e:
        #Delete files from tempSlides Folder if an error occurs!

        for temp_pptxSlide in tempNames:
            file_path = os.path.join(prsDir, temp_pptxSlide)
            try:
                if os.path.isfile(file_path) or os.path.islink(file_path):
                    os.unlink(file_path)
                elif os.path.isdir(file_path):
                    shutil.rmtree(file_path)
            except Exception as e:
                print('Failed to delete %s. Reason: %s' % (file_path, e))
        print("An error ocurred during the upload of "+filename)
        print(e)
        return fileName+"_"+".pptx"
    else:
        #Delete files from tempSlides Folder if no exception happened

        for temp_pptxSlide in tempNames:
            file_path = os.path.join(prsDir, temp_pptxSlide)
            try:
                if os.path.isfile(file_path) or os.path.islink(file_path):
                    os.unlink(file_path)
                elif os.path.isdir(file_path):
                    shutil.rmtree(file_path)
            except Exception as e:
                print('Failed to delete %s. Reason: %s' % (file_path, e))
        return None