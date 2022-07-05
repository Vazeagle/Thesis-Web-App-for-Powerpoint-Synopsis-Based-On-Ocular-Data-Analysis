from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
import pandas as pd
import numpy as np
import os
import tkinter
from classes import *
#from eye_tracking import *
import json
from PIL import ImageFont
from PIL import Image as pil_img
import base64


# Get resolution of screen
def get_screen_res():
    root = tkinter.Tk()
    screen_res = [root.winfo_screenwidth(), root.winfo_screenheight()]
    return screen_res

def get_fonts():
    #Get Fonts List
    with open('fonts.txt') as f:
        data = f.read()
    jsFonts = json.loads(data)
    return jsFonts

def get_pil_text_size(text, font_size, font_name):  #using pillow get the size and font of the specific text input
    print("font_name=",font_name)
    print("font_size=",font_size)
    print("cwd=",os.getcwd())
    #fpath = os.getcwd()+"\\fonts\\"+font_name  #optional font folder
    font = ImageFont.truetype(font_name, font_size)
    size = list(font.getsize(text))
    #return in pixels and not in Pt
    size[0] = size[0] * 1.33
    size[1] = size[1] * 1.33
    size = tuple(size)
    return size



slide_dict = {} #hold slide objects with names for better vars.

##Get Gaze Data
##Occurrences of pptx object
#def getGazeData(dataFrame, paragraphlines, xPixelStart, yPixelStart, lineHeight, LineMargin):   #paragraphlines is an array that has the width of the lines of a paragraph element inside a Textbox #AND ALSO bellow the loop happens for every par so all is ok!!!
#    occurrences = 0
#    yPixelEnd = yPixelStart - lineHeight - LineMargin
#    for lineWidth in paragraphlines:
#        xPixelEnd = xPixelStart + lineWidth
#        print("getGazeData")
#        objectData = dataFrame.loc[(dataFrame['GazeY'] <= yPixelStart) & (dataFrame['GazeY'] >= yPixelEnd) & (dataFrame['GazeX'] >= xPixelStart) & (dataFrame['GazeX'] <= xPixelEnd)]
#        print("OCC",occurrences)
#        print("DATAFRAMEC:\n",objectData)
#        occurrences = len(objectData) + occurrences
#    return occurrences
#

#   Convert emu(english metrix units) to pixels
def emu_to_pixels(emuValue: int):
    pixelValue = float(emuValue)/9525
    return pixelValue

def getScaleMultiplier(resolution):
    screen_res=get_screen_res()
    scaleMultiplier = screen_res[0]/resolution[0]
    return scaleMultiplier

def ScaleToScreenResolution(resolution,pixels):
    # The resolution for both width and height goes proportional for both
    # This function works dynamically with every slide resolution for better compatibility and functionality
    screen_res=get_screen_res()
    scaleMultiplier = screen_res[0]/resolution[0]
    if(isinstance(pixels,list)):
        result = []
        for elem in pixels:
            result.append(elem*scaleMultiplier)
        print("ORIGINAL PIXELS: ", pixels)
        print("scaleMultiplier: ", scaleMultiplier," Result: ", result,"\n")
    else:
        result = pixels*scaleMultiplier
        print("ORIGINAL PIXELS: ", pixels)
        print("scaleMultiplier: ", scaleMultiplier," Result: ", result,"\n")
    return result

def splitInnerShapes(shapeElement):
    #This function will split in virtual sub-textboxes based on the text of the textbox
    #Also this function takes some user error correcting methods to have more accurate data in case of e.g. new line with only empty spaces.
    splittedText = shapeElement.text.splitlines()
    print(splittedText)
    print("type of splittedText:",type(splittedText))
    for text in splittedText:
        if not text:
            print("null")
        if text:
            print("not null")
        if text.isspace():
            print("is space or tab", text)
            splittedText = list(map(lambda item: item.replace(text,""), splittedText))

    return splittedText

def getTextboxInfo(textboxShape):
    # reading the data from the file
    fontSize = None
    fontName = None
    #This function returns a dictionary having the font, the paragraph spacing and the margins
    for paragraph in textboxShape.text_frame.paragraphs:
        print("Paragraph font:",  paragraph.font)
        #if the following are None then they have their default value.
        print("paragraph font in PT:",  paragraph.font.size)        ############### NOT NEEDED WANT LOWER HIERARCHY
        print("paragraph font Name:",  paragraph.font.name)         ############### NOT NEEDED WANT LOWER HIERARCHY

        print("Paragraph line spacing :",  paragraph.line_spacing)  ############### NOT NEEDED WANT LOWER HIERARCHY
        print("THIS IS PRINT IN LINE 109!")
        ###---------------------------------------------------------------------------------------------------------
        ###---------------------------------------------------------------------------------------------------------
        ###---------------------------------------------------------------------------------------------------------
        ###---------------------------------------------------------------------------------------------------------
        #***************************This is added here beacuse of the subtitle element of the placeholder************************
        if paragraph.line_spacing is None: #beacause of negative values in lineSpacer in DB due to reading line_spacing as 0
            lineSpacing = 0 #default value is 1 that means no space between paragraphs, here we put 0 because it acts as a scale for yDistance variable
        else:
            lineSpacing = paragraph.line_spacing-1 #this will give the scale of increase for later use to get the pixel size
        ###---------------------------------------------------------------------------------------------------------
        ###---------------------------------------------------------------------------------------------------------
        ###---------------------------------------------------------------------------------------------------------
        #this will be used as a miltiplier, the space between paragraphs based on the font-height will be x times bigger
        access_flag=True    #NEW TO ENTER ONLY ONCE THE LINE SPACING
        for run in paragraph.runs:
            #We make a rule that only one font will be used in a paragraph or else there would be chaos
            #Even if the font exist sometimes it can't be read due to pptx limitations so we use this as a security measure
           #if run.font is not None:
            print("Run:", run)
            print("RunText:", run.text)
            if run.font.name is not None:
                print("print hello!!!!!!!@@@#@#@#")
                print("RUN: ",run)
                print("run.font.name: ", run.font.name)
                print("run.font.size: ", run.font.size)
                #Handle default font size 18 which is returned as None
                if run.font.size is None:
                    fontSize = 18   #Default font size in pt
                else:
                    fontSize = int(run.font.size.pt)
                jsFonts = get_fonts()
                fontName = jsFonts[run.font.name] # this will use the fonts dictionary to get the .tff name of the font 

                if paragraph.line_spacing is None: #beacause of negative values in lineSpacer in DB due to reading line_spacing as 0
                    lineSpacing = 0 #default value is 1 that means no space between paragraphs, here we put 0 because it acts as a scale for yDistance variable
                elif paragraph.line_spacing is not None and access_flag:
                    access_flag = False
                    lineSpacing = paragraph.line_spacing-1 #this will give the scale of increase for later use to get the pixel size
            else:
                print("No font was specifically defined in this run: ", run.text, " of this paragraph: ", paragraph.text)
                print("If a specified font is found we will use it as a default!")
                raise Exception("No font was specified for this TextBox: ", "'",textboxShape.text,"'")
            #could implement a flag to not include this paragraph on the synopsis and continue with the rest
            #if run.font.name is None:
            #    raise Exception("No font was specified for this TextBox: ", "'",textboxShape.text,"'")
    #Get TextBox Margin Properties in Pixels
    leftMargin = textboxShape.text_frame.margin_left/9525
    rightMargin = textboxShape.text_frame.margin_right/9525
    topMargin = textboxShape.text_frame.margin_top/9525
    botMargin = textboxShape.text_frame.margin_bottom/9525
    print("line 171 print, hello!")
    #Because sometimes if we leave the linespacing undefined in PowerPoint th value of linespacing could be negative because of lines 127 and 156
    if (lineSpacing < 0):
        lineSpacing = 0
    print("lineSpacing is:",lineSpacing)
    return {"fontSize": fontSize, "fontName": fontName, "lineSpacing": lineSpacing, "leftMargin": leftMargin, "rightMargin": rightMargin, "topMargin": topMargin, "botMargin": botMargin}

def locateShape(shapeElement, resolution):
    # The elements place x and y will be the top left of the text box
    xAxis = (shapeElement.left/9525) # Slide WIDTH - distance of the left edge of this shape from the left edge of the slide
    yAxis = resolution[1] - (shapeElement.top/9525) #  Slide HEIGHT - distance of the top edge of this shape from the top edge of the slide
    yAxis = 748 - (shapeElement.top/9525)
    width = shapeElement.width/9525 #Distance between left and right extents of shape
    height = shapeElement.height/9525 #Distance between top and bottom extents of shape
    #This function will produce the results in pixels based on the default-set resolution of the powerpoint slide that the shape is in.
    return {"xAxisStart": xAxis, "yAxisStart": yAxis, "width": width, "height": height}

def textbox(slideRes, slideNum, text_shape, group, objCounter, scaleMultiplier):
    shapes_results=[]
    textBoxInfo = getTextboxInfo(text_shape)
    #{"fontSize": fontSize, "fontName": fontName, "lineSpacing": lineSpacing, "leftMargin": leftMargin, "rightMargin": rightMargin, "topMargin": topMargin, "botMargin": botMargin}
                
    #height step in pixels for specific font of paragraph
    yStep = get_pil_text_size(" ",textBoxInfo["fontSize"], textBoxInfo["fontName"])[1]
    margin = (yStep*textBoxInfo["lineSpacing"])
    print("margin",margin)


    shapeLoc = locateShape(text_shape, slideRes)
    # Textboxes real shape location without margins:
    realXaxisStart = shapeLoc["xAxisStart"] + textBoxInfo["leftMargin"]
    realYaxisStart = shapeLoc["yAxisStart"] - textBoxInfo["topMargin"]# - (textBoxInfo["lineSpacing"]*yStep)
    realParWidth = shapeLoc["width"] - textBoxInfo["rightMargin"] - textBoxInfo["leftMargin"]
    realHeight = shapeLoc["height"] - textBoxInfo["botMargin"] - textBoxInfo["topMargin"]

    paragraphsList = splitInnerShapes(text_shape)
    paragraphHeightStart = realYaxisStart ############################################################ - (textBoxInfo["lineSpacing"]*yStep)#THIS WILL REMOVE THE TOP LINESPACING OF THE FIRST PARAGRAPH
    nextParYstart = 0
    space = 0
    par_counter = 0 #only counts actual paragrapghs not empty ones
                

    ####count spaces "" before and after par text STUPID WAY!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!
    ####------------------------------------------------------------------------------------------
    ####------------------------------------------------------------------------------------------
    ####---------------------------------THIS FOR LOOP IS TO GET THE CORRECT ParToParDistanceList!---------------------------------------------------------
    ParToParDistanceList = []   #This List holds the distance on y-axis between continuing paragraphs
    for par in paragraphsList:
        if not par: #If paragraph line is empty and has no text
            space +=1
        else:
            paragraphLines = 1 #initialize paragraph lines counter
            realParTextLength = get_pil_text_size(par,textBoxInfo["fontSize"], textBoxInfo["fontName"])[0] # This is the paragraph text length based on it's font
            pptxParlength = realParTextLength + textBoxInfo["leftMargin"] + textBoxInfo["rightMargin"] # This is the real length because we have to consider the margins of the textBox
            if pptxParlength > shapeLoc["width"]:
                if (pptxParlength/shapeLoc["width"]) > int(pptxParlength/shapeLoc["width"]):
                    paragraphLines = int(pptxParlength/shapeLoc["width"])+1
            #yDistance is the between distance of two paragraphs in Pixels in the yAxis
            #old version was ParToParDistance = space*yStep + (space+1)*(yStep*textBoxInfo["lineSpacing"])
            ParToParDistance = (space*yStep) + ((space*2)*(yStep*textBoxInfo["lineSpacing"])) #This will contain the space height based on the line and the two line spacings(top and bot) of the empty line/lines!!
            if ParToParDistance == 0:#if ParToParDistance is 0 because space=0 it meants that the paragraphs are one after an other
                print("LMAO2",space)
                ParToParDistance = 2 * (yStep*textBoxInfo["lineSpacing"]) # "1*" would only put us to the top side of the next line line spacing if we put "2*" we will have the text
            else:
                print("fuckyouSpacing2",space)
                print("ystep",yStep)
                print("linespacing",textBoxInfo["lineSpacing"])
                ParToParDistance = ParToParDistance + 2*(yStep*textBoxInfo["lineSpacing"]) #ADD bot linespace of previous line and also the top linespacing(of next line) to get the height where the next line text is if before there was a space line\lines
                print("ParToParDistanceJUSTWTIHSPACES!",ParToParDistance - 2*(yStep*textBoxInfo["lineSpacing"]))
                print("CorrectParToParDistance",ParToParDistance)
            space = 0 #re-initialize space between paragraphs
            ParToParDistanceList.append(ParToParDistance)
    ParToParDistanceList.append(0)#This value is needed for the last paragraph in order to be in bounds!!!!!!!
    ####------------------------------------------------------------------------------------------
    ####------------------------------------------------------------------------------------------
    ####------------------------------------------------------------------------------------------
    ####------------------------------------------------------------------------------------------
                
    for par in paragraphsList:
        parLinesWidth = []
        if not par: #If paragraph line is empty and has no text
            space +=1
            print("Paragraph line is empty\n")
            #We need the height of the font in get it with space?
        else:
            print("Paragraph line is full")
            par_counter+=1
            paragraphLines = 1 #initialize paragraph lines counter
            realParTextLength = get_pil_text_size(par,textBoxInfo["fontSize"], textBoxInfo["fontName"])[0]
            pptxParlength = realParTextLength + textBoxInfo["leftMargin"] + textBoxInfo["rightMargin"] # This is the real length because we have to consider the margins of the textBox
            if pptxParlength > shapeLoc["width"]:
                if (pptxParlength/shapeLoc["width"]) > int(pptxParlength/shapeLoc["width"]):
                    paragraphLines = int(pptxParlength/shapeLoc["width"])+1
            i=1 #i is like a line counter inside the paragraph
            #parLinesWidth = []
            while i <= paragraphLines:
                if i == paragraphLines:
                    print("realParTextLength",realParTextLength)
                    print("realParWidth * i",realParWidth * i)
                    print("realParWidth",realParWidth)
                    print("i",i)
                    parLinesWidth.append(realParTextLength - (realParWidth * (i-1)))
                else:
                    parLinesWidth.append(realParWidth)
                i+=1
            print("PARA:",parLinesWidth)

            #Old version before malakia1821 was
            #  if par_counter==1 and space > 1:            
            if par_counter==1 and space >= 1: #This will handle the use case where before the first par there are empty lines- empty paragraphs FILLED WITH SPACES?
                #MARGIN IN PIXELS margin =  yStep*textBoxInfo["lineSpacing"]
                #OLD VERSION OF PARHEIGHT parHeight = paragraphLines*yStep + (paragraphLines-1)*(yStep*textBoxInfo["lineSpacing"]) #This is the height of the empty paragraphs + the line spacing between them(we don't include the last spacing where it is connected with a par that has text below because we count that later on)
                parHeight = paragraphLines*yStep + (paragraphLines*2)*(yStep*textBoxInfo["lineSpacing"])
                #yDistance is the between distance of two paragraphs in Pixels in the yAxis
                #old version was: ParToParDistance = space*yStep + (space+2)*(yStep*textBoxInfo["lineSpacing"]) #THIS IS space+2 because we have deleted the initial space from realYaxisStart
                ParToParDistance =ParToParDistanceList[par_counter-1] - (yStep*textBoxInfo["lineSpacing"])# minus this linespacing cause for the first par with spaces before it we are putting one extra on ParToParDistanceList Creation!!
                #SOS SOS SOS SOS ON ParToParDistanceList we sum the whole spacing of the spaces before and after the line
                # but in  paragraphHeightStart which is initialized as "realYaxisStart" we have removed the tope line spicaing so this action
                # paragraphHeightStart - ParToParDistance should get us lower than we want 
                # but NO THIS IS CORRECT CAUSE THE CURREN PAR THAT HAS TEXT HAS ALSO THE TOP LINE SPACING ACTIVE SO IT WILL CANCEL OUT WITH THE PREVIOUS  SUBSTRACTION 
                # MAYBE IT WOULD BE MORE CORECT IF I WROTE THE SUBSTRACTION HERE!!!
                # ALSO IN THIS WAY THE nextParYstart NEEDS FIXING!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!
                # THE FIX IS TO FIX parHeight to
                # think if for parHeight we need the lower linespacing or not!!!!!!!!!!!!!!!!@@@@@@@@@@@@########$$$$$$$$$$$$%%%%%%%%%%%%%%
                #parHeight = paragraphLines*yStep + (paragraphLines*2)*(yStep*textBoxInfo["lineSpacing"])
                ##########OMOIA FIX KAI THN ELSE APO KATW!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!
                paragraphHeightStart = paragraphHeightStart - ParToParDistance # stores first paragraph beginning
                nextParYstart = parHeight - 2*(yStep*textBoxInfo["lineSpacing"])+ ParToParDistanceList[par_counter] # we remove the bot of the linespacing of the current line beacause it's already sumed in ParToParDistanceList
                print("malakia182156")
                print("paragraphHeightStart:", paragraphHeightStart)
                print("ParToParDistance:", ParToParDistance)
                print("parHeight:", parHeight)
                print("yStep:", yStep)
                print("marginSpace:", (space+1)*(yStep*textBoxInfo["lineSpacing"]))
                print("nextParYstart:",nextParYstart)
                space = 0 #re-initialize space between paragraphs
            else:
                #MARGIN IN PIXELS margin =  yStep*textBoxInfo["lineSpacing"]
                #old version parHeight parHeight = paragraphLines*yStep + (paragraphLines-1)*(yStep*textBoxInfo["lineSpacing"])
                parHeight = paragraphLines*yStep + (paragraphLines*2)*(yStep*textBoxInfo["lineSpacing"])
                #yDistance is the between distance of two paragraphs in Pixels in the yAxis
                if par_counter==1:  #Handling the first par.
                    #in case that there is no empty lines before the first par we have to remove one extra line spacing (note we could mod-change the ParToParDistanceList for this instead of doing it here!)
                    paragraphHeightStart = paragraphHeightStart - (yStep*textBoxInfo["lineSpacing"])
                else:
                    paragraphHeightStart = paragraphHeightStart - nextParYstart # stores each paragraph height start
                ParToParDistance = ParToParDistanceList[par_counter]# this is the par to par distance that was calculated before!
                print("malakia12")
                print("ParToParDistanceList",ParToParDistanceList)
                print("paragraphHeightStart:", paragraphHeightStart)
                print("nextParYstart:", nextParYstart)
                print("ParToParDistance:", ParToParDistance)
                print("parHeight:", parHeight)
                print("yStep:", yStep)
                print("marginSpace:", (space+1)*(yStep*textBoxInfo["lineSpacing"]))
                nextParYstart = parHeight - 2*(yStep*textBoxInfo["lineSpacing"]) + ParToParDistance #calculates next paragraphs height start on y axis in pixels
                print("nextParYstartForNextLINE10:", nextParYstart)
                space = 0 #re-initialize space between paragraphs
                        

            #change to correct resolution
            print("test")
            testingS = ScaleToScreenResolution(slideRes, ParToParDistanceList)
            print("test")
            realHeight = parHeight-2*(yStep*textBoxInfo["lineSpacing"])
            print("\nrealHeight before Conversion to Scale:",realHeight,"\n")

            realXaxisStartFinal = ScaleToScreenResolution(slideRes, realXaxisStart)
            realYaxisStart = ScaleToScreenResolution(slideRes, paragraphHeightStart)
            realParWidthPixel = ScaleToScreenResolution(slideRes, realParWidth)
            realHeight = ScaleToScreenResolution(slideRes, realHeight)
            #realHeight is the height that also contains the lineSpacing from both up and down the row of a text
            yPixelStep = ScaleToScreenResolution(slideRes, yStep)#yPixelStep is not yStep because it will call it again and again and multiply it by scale var and deform the data that use ystep above
            Pixelmargin = ScaleToScreenResolution(slideRes, margin)#same as above for the margin!
            parLinesWidth = ScaleToScreenResolution(slideRes,parLinesWidth)
            print("HEIGHT WITH NO LINESPACING",parHeight-2*(yStep*textBoxInfo["lineSpacing"]))

                        
                        
            #Get Occurrences- Duration of Focus
            print(parLinesWidth)
            OCnum = 0 #getGazeData(df, parLinesWidth, realXaxisStart, realYaxisStart, yPixelStep, Pixelmargin)
            print("RESULT OF getGazeData:",OCnum)

            #increase object counter  of the slide by one
            objCounter=objCounter+1

            shapeInfo = TextBox(realXaxisStartFinal, realYaxisStart, realParWidthPixel, realHeight, "TEXT_BOX", slideNum, objCounter, par, parLinesWidth, Pixelmargin, yPixelStep, textBoxInfo["fontName"], textBoxInfo["fontSize"], textBoxInfo["leftMargin"], textBoxInfo["rightMargin"], textBoxInfo["topMargin"], textBoxInfo["botMargin"], scaleMultiplier, group)#  Create object of class TextBox
            #shfape_dict['Shape_' + str(shape_counter) + '_slide_' + str(slideNumber) + '_par_' + str(par_counter)] = shapeInfo# Insert Object to dictionary with key Slide_'#slide_number'
            #SOS SOS SOS SOS DO I NEED TO SAVE The shape to dict with dynamic names because every with only one shapeInfo we will only save there the last one on memory
            #ALSO SOS SOS SOS SOS WRITE TO REPORT THAT IT VIEWS EAXH PARAGRAPH AS A SEPARATE TEXTBOX ELEMENT EVEN IF IT IS IN ONE TEXTBOX.
            shapes_results.append(shapeInfo.getValues())
            

            #To Do:
            #SOS prepei na exoume kai use case stin periptwsi pou bazei to text grouped eite sthn mesh eite sta dexia anti gia to default aristera  
            #SOS OTAN EXOUME TWRA ESTW TITLO ME TOPOTHETISI STHN MESH OPWS EINAI O KWDIKAW TWRA THA KANEI STA MAP TO WIDTH APO THN ARISTERH PLEURA KAI OXI STHN MESH TO KEIMENO OPOTEW PROBLEM!!!     
            #Complete the function to get data from the dataframe using x and y and get the timestamps and occurences from the eyetracking.py
            #Open subproccess for the tobiistream.exe and create a interupt to stop it so it can continue later on  with a flag for the next slide.
            #Code Refinement and remove unessesary code from classes.py
    return shapes_results, objCounter

def picture(pptxName, slideRes, slideNum, picture_shape, group, objCounter, scaleMultiplier):

    image_loc = locateShape(picture_shape, slideRes)
    print("This is an image:",image_loc)

    width = [image_loc["width"]]
    OCnum = 0 #getGazeData(df, width, image_loc["xAxisStart"], image_loc["yAxisStart"], image_loc["height"], 0)
    print("RESULT OF getGazeData for Picture:",OCnum)

    #increase object counter  of the slide by one
    objCounter=objCounter+1

    img_blob = picture_shape.image.blob
    img_extensionType = picture_shape.image.ext
    blob_save_loc = os.getcwd() + "\\uploads\\" + pptxName.rsplit('.', 1)[0] + "\\" + str(slideNum) + "_" + str(objCounter) + ".blob"
    img_save_loc = os.getcwd() + "\\uploads\\" + pptxName.rsplit('.', 1)[0] + "\\" + str(slideNum) + "_" + str(objCounter) + "." + img_extensionType #remove the .pptx extension of the file for the correct filename
    img_loc = "\\uploads\\" + pptxName.rsplit('.', 1)[0] + "\\" + str(slideNum) + "_" + str(objCounter) + "." + img_extensionType
    #Image.open(base64.decodebytes(img_blob))
    with open(img_save_loc, "wb") as blob_save_loc:
        blob_save_loc.write(img_blob)#SAVE BLOB CAUSE IF WE CALL IT TO DECODE TI AND SAVE IT AS AN PICTURE EXTENSION IT FAILS!!!!
        #pic_to_save.write(base64.decodebytes(img_blob)) #FAILS!
        img = pil_img.open(img_save_loc)
        #img.show()
        img.save(img_save_loc)#save image to specified path!

    #FIX PIXEL VALUES TO BE INSERTED BASED ON THE SCREEN RESOLUTION AND NOT THE SLIDE RESOLUTION!
    realxAxisStart = ScaleToScreenResolution(slideRes, image_loc["xAxisStart"])
    realyAxisStart = ScaleToScreenResolution(slideRes, image_loc["yAxisStart"])
    realWidth = ScaleToScreenResolution(slideRes, image_loc["width"])
    realHeight = ScaleToScreenResolution(slideRes, image_loc["height"])


    imageInfo = Image(realxAxisStart, realyAxisStart, realWidth, realHeight, "PICTURE", slideNum, objCounter, img_loc, img_extensionType, scaleMultiplier, group)
    return imageInfo.getValues(), objCounter
    
    #Check also https://stackoverflow.com/questions/58804348/convert-blob-image-to-png-jpg-with-python
    #CHECK https://stackoverflow.com/questions/62864082/copy-slide-with-images-python-pptx AND SOS LOOK LAST COMMENT ABOUT UNIQUINESS
    # image_loc = {"xAxisStart": xAxis, "yAxisStart": yAxis, "width": width, "height": height}
    # I could add image rotation here????????????????

def gatherData(powerpointPath, powerpointName):

    prs = Presentation(powerpointPath)
    textbox_dict_list=[]
    picture_dict_list=[]
    slide_dict_list=[]
    for slide in prs.slides:
        objectCounter=0 # initialise counter for counting objects in for each slide it increases for each sub object we count- create form par images groups etc!
        width = prs.slide_width/9525
        height = prs.slide_height/9525
        slideResolution = [width, height]
        slideNumber = prs.slides.index(slide)
        scaleMultiplier = getScaleMultiplier(slideResolution) # slide scaleMultiplier

        slideInfo = Slide(width, height, slideNumber)#  Create object of class Slide
        slide_dict_list.append(slideInfo.getValues())
        #slide_dict['Slide_' + str(prs.slides.index(slide))] = slideInfo# Insert Object to dictionary with key Slide_'#slide_number'

        shape_counter = 0
        group_counter = 0
        for shape in slide.shapes:
            print("Type of shape", shape.shape_type)#   Type of shape!
            shape_counter+=1    #groups are seen as one shape
            #For Textboxes
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                textboxInfo = textbox(slideResolution, slideNumber, shape, None, objectCounter, scaleMultiplier)#None because of no associated group
                for dict in textboxInfo[0]: #because the textbox function returns a list that contains dictionaries full with information from class objects methods
                    textbox_dict_list.append(dict)  #append dictionary with element info
                objectCounter = textboxInfo[1]  #return element counter to increase correctly the counter
                    

            #print("TEXTBOXES INFO\n")
            #print(shape_dict)
            
            #IF SHAPE IS A PICTURE
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pictureInfo = picture(powerpointName, slideResolution, slideNumber, shape, None, objectCounter, scaleMultiplier)#None because of no associated group
                picture_dict_list.append(pictureInfo[0]) #append dictionary with element info
                objectCounter = pictureInfo[1] #return element counter to increase correctly the counter
                print("This is an image")

            #IF SHAPE IS A GROUP OF OTHER SHAPES(TEXT_BOX AND PICTURE)
            #!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!
            #!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!
            #!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!

            #SOS WE NEED TO CONNECT THE PICTURES WITH THE TEXT WITH SOME WAY MAYBE A SPECIFIC UNIQUE KEY
            #I COULD CREATE A STORED PROCEDURE OR TRIGGER IN MYSQL DATABASE THAT WILL HANDLE THAT DATA.
            #MAYBE DO THIS
            #CREATE A NEW COLUMN IN THE TABLE OF THE SHAPES CALLED "CORELLATION_ID"
            #AUTO TO ID THA DEIXNEI PALI PISW STO ID TOU OBJECT(EIKONA H KEIMENO POY EISIXTH PRWTO)
            #OPOTE SE KA8E MIA EISAGWGH SUB-SHAPE TO EPOMENO PRAGMA EINAI CALL STHN DATABASE GIA TO ID TOU WSTE AUTO NA PAEI SAN CORELLATION ID STA EPOMENA

            #!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!
            #!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!
            #!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                print("This is a group")
                groupName = group_counter
                group_counter+=1
                #this will work for simple 1st level grouping it wont iterate in the more inside groups!!!
                #if we want to iterate more check the bellow link!!!!
                for sub_shape in shape.shapes:
                    if sub_shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX: 
                        print("This is a TextBox")
                        textboxInfo = textbox(slideResolution, slideNumber, sub_shape, groupName, objectCounter, scaleMultiplier)
                        for dict in textboxInfo[0]:#because the textbox function returns a list that contains dictionaries full with information from class objects methods
                            textbox_dict_list.append(dict)#append dictionary with element info
                        objectCounter = textboxInfo[1] #return element counter to increase correctly the counter
                    if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        print("This is an Image")
                        pictureInfo = picture(powerpointName, slideResolution, slideNumber, sub_shape, groupName, objectCounter, scaleMultiplier)
                        picture_dict_list.append(pictureInfo[0])#append dictionary with element info
                        objectCounter = pictureInfo[1] #return element counter to increase correctly the counter
            # https://stackoverflow.com/questions/51701626/how-to-extract-text-from-a-text-shape-within-a-group-shape-in-powerpoint-using

            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                print("This is a PLACEHOLDER")
                phf = shape.placeholder_format
                print('%d, %s' % (phf.idx, phf.type))
                #this will work for simple 1st level grouping it wont iterate in the more inside groups!!!
                #if we want to iterate more check the bellow link!!!!
                if phf.type == PP_PLACEHOLDER.TITLE or phf.type == PP_PLACEHOLDER.SUBTITLE or phf.type == PP_PLACEHOLDER.CENTER_TITLE or phf.type == PP_PLACEHOLDER.VERTICAL_TITLE or phf.type == PP_PLACEHOLDER.BODY: 
                    print("This is a Text inside a placeholder")
                    print("TEXT :",shape.text)
                    if (shape.text !=""):#if user hasn't filled all placeholders with text and has left some on default it will raise an error on getting text info as width length etc. so we have to use this if to check if there is text inside it!
                        textboxInfo = textbox(slideResolution, slideNumber, shape, None, objectCounter, scaleMultiplier)
                        for dict in textboxInfo[0]:#because the textbox function returns a list that contains dictionaries full with information from class objects methods
                            textbox_dict_list.append(dict)#append dictionary with element info
                        objectCounter = textboxInfo[1] #return element counter to increase correctly the counter
                if phf.type == PP_PLACEHOLDER.PICTURE:
                    #https://stackoverflow.com/questions/64903569/attributeerror-slideplaceholder-object-has-no-attribute-insert-picture
                    #HAS TO BE PICTURE-PLACEHOLDER IN ORDER TO WORK CORRECTLY!!
                    print("This is an Image inside a placeholder")
                    pictureInfo = picture(powerpointName, slideResolution, slideNumber, shape, None, objectCounter, scaleMultiplier)
                    picture_dict_list.append(pictureInfo[0])#append dictionary with element info
                    objectCounter = pictureInfo[1] #return element counter to increase correctly the counter
            
    return textbox_dict_list, picture_dict_list, slide_dict_list



#test dictionary with slide classes.
#print("slide_dict: ",slide_dict)

#test = slide_dict.values()
#for class_elem in test:
#    class_elem.getSlideNum()



#TO DOOOOOOOOOOOOO

#check ParToParDistance it saves wrong the data maybe use a list? or reverse the current par list!!!!

#ABOUT TIMESTAMP MAYBE A FUNCTION TO SEE WHEN THE USER EYES GET OUT OF THE X,Y COORDINATES OF THE OBJECT AND THEN USE THIS TIMESTAMP AS MAX,
#THEN CHECK WHEN AND IF HE REENTERS THIS OBJECT AND CRAETE A NEW MIN-MAX AFTER HE SEES OUTSIDE SO WHEN HE STOPS THE DIFF MAX-MIN WILL BE ADDED TO THE FINAL TIME TO SEE
#HOW MUCH TIME THE USER HAS FOCUSED ON AN OBJECT