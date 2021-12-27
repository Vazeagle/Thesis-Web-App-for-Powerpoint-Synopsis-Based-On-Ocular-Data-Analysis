from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor

import pandas as pd
import numpy as np
import os
import tkinter


#	Get resolution of screen
root = tkinter.Tk()
screen_width = root.winfo_screenwidth()
screen_height = root.winfo_screenheight()

#	Get current working dir
cwd = os.getcwd()

#os.system(TobiiStream)
#POWERPOINT PPTX--------------------------------------------------
prs = Presentation('test2.pptx')

#layout = prs.slide_masters[0].slide_layouts[0]
#slide = prs.slides.add_slide(layout)
# 
#text_box = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(5), Cm(2))
#text_box.text = "My slide"
# 
#text_box.fill.solid()
#text_box.fill.fore_color.rgb = RGBColor(0, 0, 255)
#
#prs.save('test_python.pptx')
#print("Saved it!!!")

# specifix text elements of slides
slide_index=0
slide_text_shapes_num=0
slide_text_num_info_dict={}
slide_text_data_info_dict={}
slide_data_dict={}

for slide in prs.slides:
    print("Slide number: ",prs.slides.index(slide))
    print("SLIDE WIDTH: ",prs.slide_width/9525)
    print("SLIDE HEIGHT: ",prs.slide_height/9525)
    for shape in slide.shapes:
        print("Type of shape", shape.shape_type)#   Type of shape!
        if hasattr(shape, "text"):
            print("Margins of TextBox:")
            print("Left Margin: ", shape.text_frame.margin_left/9525)
            print("Right Margin: ", shape.text_frame.margin_right/9525)
            print("Top Margin: ", shape.text_frame.margin_top/9525)
            print("Bottom Margin: ", shape.text_frame.margin_bottom/9525)
            print("\n")
            #print("Paragraph:", shape.text_frame.paragraphs[0])
            #shape.text_frame.paragraphs[0].font.size = Pt(30)
            for paragraph in shape.text_frame.paragraphs:
                print("paragraph font in PT:",  paragraph.font.size)
                print("paragraph font Name:",  paragraph.font.name)
                print("Paragraph line spacing :",  paragraph.line_spacing)
                for run in paragraph.runs:
                    print("Run:", run)
                    #run.font.size = Pt(25) do everything 25 in pt font
                    #print("\nFont size:", run.font.size)

                    #Notes
                    #NEED TO CREATE A TRY CATCH IF USER HADN'T DEFINED Specific Run Font
                    #Add it to specific elem class as NoneType.
                    #SOS
                    #CREATE Function to split pptx objects mostly textboxes paragraphs with \n
                    #and function to count occurenceses of elements based on x,y data of eyetracker and timestamp min and max for the object and create the map for the elements-objects of pptx.

                    #print("\nFont sizePt:", (run.font.size/9525)/1.333)
                    #print("\nFont sizePt:", run.font.size.pt)
                    #print("\nFont Name:", run.font.name)

                    #print("\nFont sizePixels:", (run.font.size/9525)*1.333)
            #for paragraph in shape.text_frame.paragraphs:
            #    paragraph.font.size = Pt(30)
            #shape.rotation= 45.0 #works!!!
            #print("Font:", shape.text_frame.font)
            print(shape.text)
            print(shape.name)
            print("left: ",shape.left/9525)# ΑΠΟΣΤΑΣΗ ΤΟΥ ΤΕΤΡΑΓΩΝΟΥ ΤΟΥ SHAPE ΑΠΟ ΤΟ ΑΡΙΣΤΕΡΟ ΑΚΡΟ Read/write. Integer distance of the left edge of this shape from the left edge of the slide, in English Metric Units (EMU)
            print("top: ",shape.top/9525) # ΑΠΟΣΤΑΣΗ ΑΠΟ ΤΟ ΑΝΩ ΣΗΜΕΙΟ ΤΟΥ ΣΧΗΜΑΤΟΣ ΜΕΧΡΙ ΤΗΝ ΚΟΡΥΦΗ ΤΗΣ ΣΕΛΙΔΑΣ Read/write. Integer distance of the top edge of this shape from the top edge of the slide, in English Metric Units (EMU)
            print("width: ",shape.width/9525) # ΜΕΓΕΘΟΣ ΣΧΗΜΑΤΟΣ ΣΕ MHKOS????Read/write. Integer distance between left and right extents of shape in EMUs
            print("height: ",shape.height/9525)# ΜΕΓΕΘΟΣ ΣΧΗΜΑΤΟΣ ΣΕ ΥΨΟΣ Read/write. Integer distance between top and bottom extents of shape in EMUs
            print(shape.element)
            #print(shape.text.size)
            slide_text_data_info_dict[slide_text_shapes_num] = shape.text
            slide_text_shapes_num+=1
    slide_data_dict[slide_index] = slide_text_data_info_dict.copy() #need the copy or else when I clear the dict it will dynamically clear this too!
    slide_text_data_info_dict.clear()
    slide_text_num_info_dict[slide_index] = slide_text_shapes_num
    slide_index+=1
    slide_text_shapes_num = 0 #need to re initialise for next slide
prs.save('test2.pptx')
print("slide_text_num_info_dict: \n",slide_text_num_info_dict)
print("slide_data_dict: \n",slide_data_dict)

#for slide in prs.slides:
#    for shape in slide.shapes:
#        print(shape.begin_x)

#for slide in prs.slides:
#    for shape in slide.placeholders:
#        print('%d %d %s' % (prs.slides.index(slide), shape.placeholder_format.idx, shape.name))
#-----------------------------------------------------------------

#   Convert emu(english metrix units) to pixels
def emu_to_pixels(emuValue: int):
    pixelValue = float(emuValue)/9525
    return pixelValue



import ctypes

def GetTextDimensions(text, points, font):
    class SIZE(ctypes.Structure):
        _fields_ = [("cx", ctypes.c_long), ("cy", ctypes.c_long)]

    hdc = ctypes.windll.user32.GetDC(0)
    hfont = ctypes.windll.gdi32.CreateFontA(points, 0, 0, 0, 0, 0, 0, 0, 0, 0, 0, 0, 0, font)
    hfont_old = ctypes.windll.gdi32.SelectObject(hdc, hfont)

    size = SIZE(0, 0)
    ctypes.windll.gdi32.GetTextExtentPoint32A(hdc, text, len(text), ctypes.byref(size))

    ctypes.windll.gdi32.SelectObject(hdc, hfont_old)
    ctypes.windll.gdi32.DeleteObject(hfont)

    return (size.cx, size.cy)

print(GetTextDimensions("This is a paragraph test Par to check if the length is represented correctly as it should ………..", 24, "Times New Roman"))
print(GetTextDimensions("This is a paragraph test Par5 to test what we see in text when the sentence continues below.", 25, "Arial"))

from PIL import ImageFont

def get_pil_text_size(text, font_size, font_name):
    font = ImageFont.truetype(font_name, font_size)
    size = font.getsize(text)
    
    return size
print(get_pil_text_size('This is a paragraph test Par to check if the length is represented correctly as it shoulddddddddddddddddd. ', 20, 'arial.ttf')[0]*1.33)
print(get_pil_text_size('This is a paragraph test Par5 to test what we see in text when the sentence continues below.', 24, 'arial.ttf')[0]*1.33)


print("\nCheck now:")
print(get_pil_text_size('Full length of powerpoint presentation check to see the pixel length accuracy of a font test test test test test test test test test test', 18, 'calibri.ttf')[0]*1.33)
print(get_pil_text_size('Full length of powerpoint presentation check to see the pixel length accuracy of a font test test test test test test test test test test.', 18, 'calibri.ttf')[0]*1.33)
print(get_pil_text_size('Test to see middle pixel length of pillow and get the values ', 18, 'arial.ttf')[0]*1.33)
print(get_pil_text_size('Second check the accuracy of counting pixel cnt', 24, 'times.ttf')[0]*1.33)

#ctypes is only for windows and its measurements arent precise
#pillow method is more precise also the values returned are in PT SO TO CONVERT TO PIXELS MULTIPLY BY 1.333
#SOS PREPEI STIS PANW METRHSEIS NA FORTWSW KAI TA ANTYISTOIXA MARGINS GIA LEFT RIGHT TOP KAI BOTTOM KAI NA YPOLOGISW ME BASH AYTA TA PERITHORIA