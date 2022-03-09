from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd
import numpy as np
import os
import tkinter
from classes import *
#from eye_tracking import *
import json
from PIL import ImageFont


# Get resolution of screen
def get_screen_res():
    root = tkinter.Tk()
    screen_res = [root.winfo_screenwidth(), root.winfo_screenheight()]
    return screen_res

def get_fonts():
    #Get Fonts List
    with open('fonts.txt') as f:
        data = f.read()
    jsFonts = json.loads(data)
    return jsFonts

def get_pil_text_size(text, font_size, font_name):  #using pillow get the size and font of the specific text input
    font = ImageFont.truetype(font_name, font_size)
    size = list(font.getsize(text))
    #return in pixels and not in Pt
    size[0] = size[0] * 1.33
    size[1] = size[1] * 1.33
    size = tuple(size)
    return size



slide_dict = {} #hold slide objects with names for better vars.

shape_dict = {} #hold shape objects from textpox or image classes.

temp_shapes_dict = {} #hold temporarily shape objects from textpox or image classes .

#Get Gaze Data
#Occurrences of pptx object
def getGazeData(dataFrame, paragraphlines, xPixelStart, yPixelStart, lineHeight, LineMargin):   #paragraphlines is an array that has the width of the lines of a paragraph element inside a Textbox #AND ALSO bellow the loop happens for every par so all is ok!!!
    occurrences = 0
    yPixelEnd = yPixelStart - lineHeight - LineMargin
    for lineWidth in paragraphlines:
        xPixelEnd = xPixelStart + lineWidth
        print("getGazeData")
        objectData = dataFrame.loc[(dataFrame['GazeY'] <= yPixelStart) & (dataFrame['GazeY'] >= yPixelEnd) & (dataFrame['GazeX'] >= xPixelStart) & (dataFrame['GazeX'] <= xPixelEnd)]
        print("OCC",occurrences)
        print("DATAFRAMEC:\n",objectData)
        occurrences = len(objectData) + occurrences
    return occurrences


#   Convert emu(english metrix units) to pixels
def emu_to_pixels(emuValue: int):
    pixelValue = float(emuValue)/9525
    return pixelValue

def ScaleToScreenResolution(resolution,pixels):
    # The resolution for both width and height goes proportional for both
    # This function works dynamically with every slide resolution for better compatibility and functionality
    screen_res=get_screen_res()
    scaleMultiplier = screen_res[0]/resolution[0]
    if(isinstance(pixels,list)):
        result = []
        for elem in pixels:
            result.append(elem*scaleMultiplier)
        print("ORIGINAL PIXELS: ", pixels)
        print("scaleMultiplier: ", scaleMultiplier," Result: ", result,"\n")
    else:
        result = pixels*scaleMultiplier
        print("ORIGINAL PIXELS: ", pixels)
        print("scaleMultiplier: ", scaleMultiplier," Result: ", result,"\n")
    return result

def splitInnerShapes(shapeElement):
    #This function will split in virtual sub-textboxes based on the text of the textbox
    #Also this function takes some user error correcting methods to have more accurate data in case of e.g. new line with only empty spaces.
    splittedText = shapeElement.text.splitlines()
    print(splittedText)
    print("type of splittedText:",type(splittedText))
    for text in splittedText:
        if not text:
            print("null")
        if text:
            print("not null")
        if text.isspace():
            print("is space or tab", text)
            splittedText = list(map(lambda item: item.replace(text,""), splittedText))

    return splittedText

def getTextboxInfo(textboxShape):
    # reading the data from the file
    fontSize = None
    fontName = None
    #This function returns a dictionary having the font, the paragraph spacing and the margins
    for paragraph in textboxShape.text_frame.paragraphs:
        print("Paragraph font:",  paragraph.font)
        #if the following are None then they have their default value.
        print("paragraph font in PT:",  paragraph.font.size)        ############### NOT NEEDED WANT LOWER HIERARCHY
        print("paragraph font Name:",  paragraph.font.name)         ############### NOT NEEDED WANT LOWER HIERARCHY

        print("Paragraph line spacing :",  paragraph.line_spacing)  ############### NOT NEEDED WANT LOWER HIERARCHY
        #this will be used as a miltiplier, the space between paragraphs based on the font-height will be x times bigger
        for run in paragraph.runs:
            #We make a rule that only one font will be used in a paragraph or else there would be chaos
            #Even if the font exist sometimes it can't be read due to pptx limitations so we use this as a security measure
           #if run.font is not None:
            print("Run:", run)
            print("RunText:", run.text)
            if run.font.name is not None:
                fontSize = int(run.font.size.pt)
                jsFonts = get_fonts()
                fontName = jsFonts[run.font.name] # this will use the fonts dictionary to get the .tff name of the font 

                if paragraph.line_spacing is None:
                    lineSpacing = 0 #default value is 1 that means no space between paragraphs, here we put 0 because it acts as a scale for yDistance variable
                else:
                    lineSpacing = paragraph.line_spacing-1 #this will give the scale of increase for later use to get the pixel size
            else:
                print("No font was specifically defined in this run: ", run.text, " of this paragraph: ", paragraph.text)
                print("If a specified font is found we will use it as a default!")
            #could implement a flag to not include this paragraph on the synopsis and continue with the rest
    if run.font.name is None:
        raise Exception("No font was specified for this TextBox: ", "'",textboxShape.text,"'")
    #Get TextBox Margin Properties in Pixels
    leftMargin = textboxShape.text_frame.margin_left/9525
    rightMargin = textboxShape.text_frame.margin_right/9525
    topMargin = textboxShape.text_frame.margin_top/9525
    botMargin = textboxShape.text_frame.margin_bottom/9525
    return {"fontSize": fontSize, "fontName": fontName, "lineSpacing": lineSpacing, "leftMargin": leftMargin, "rightMargin": rightMargin, "topMargin": topMargin, "botMargin": botMargin}

def locateShape(shapeElement, resolution):
    # The elements place x and y will be the top left of the text box
    xAxis = (shapeElement.left/9525) # Slide WIDTH - distance of the left edge of this shape from the left edge of the slide
    yAxis = resolution[1] - (shapeElement.top/9525) #  Slide HEIGHT - distance of the top edge of this shape from the top edge of the slide
    width = shapeElement.width/9525 #Distance between left and right extents of shape
    height = shapeElement.height/9525 #Distance between top and bottom extents of shape
    #This function will produce the results in pixels based on the default-set resolution of the powerpoint slide that the shape is in.
    return {"xAxisStart": xAxis, "yAxisStart": yAxis, "width": width, "height": height}

def textbox(slideRes, slideNum, text_shape, df):
    textBoxInfo = getTextboxInfo(text_shape)
    #{"fontSize": fontSize, "fontName": fontName, "lineSpacing": lineSpacing, "leftMargin": leftMargin, "rightMargin": rightMargin, "topMargin": topMargin, "botMargin": botMargin}
                
    #height step in pixels for specific font of paragraph
    yStep = get_pil_text_size(" ",textBoxInfo["fontSize"], textBoxInfo["fontName"])[1]
    margin = (yStep*textBoxInfo["lineSpacing"])


    shapeLoc = locateShape(text_shape, slideRes)
    # Textboxes real shape location without margins:
    realXaxisStart = shapeLoc["xAxisStart"] + textBoxInfo["leftMargin"]
    realYaxisStart = shapeLoc["yAxisStart"] - textBoxInfo["topMargin"] - (textBoxInfo["lineSpacing"]*yStep)
    realParWidth = shapeLoc["width"] - textBoxInfo["rightMargin"] - textBoxInfo["leftMargin"]
    realHeight = shapeLoc["height"] - textBoxInfo["botMargin"] - textBoxInfo["topMargin"]

    paragraphsList = splitInnerShapes(text_shape)
    paragraphHeightStart = realYaxisStart
    nextParYstart = 0
    space = 0
    par_counter = 0 #only counts actual paragrapghs not empty ones
                

    ####count spaces "" before and after par text STUPID WAY!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!
    ####------------------------------------------------------------------------------------------
    ####------------------------------------------------------------------------------------------
    ####------------------------------------------------------------------------------------------
    ParToParDistanceList = []
    for par in paragraphsList:
        if not par: #If paragraph line is empty and has no text
            space +=1
        else:
            paragraphLines = 1 #initialize paragraph lines counter
            realParTextLength = get_pil_text_size(par,textBoxInfo["fontSize"], textBoxInfo["fontName"])[0]
            pptxParlength = realParTextLength + textBoxInfo["leftMargin"] + textBoxInfo["rightMargin"] # This is the real length because we have to consider the margins of the textBox
            if pptxParlength > shapeLoc["width"]:
                if (pptxParlength/shapeLoc["width"]) > int(pptxParlength/shapeLoc["width"]):
                    paragraphLines = int(pptxParlength/shapeLoc["width"])+1
            #yDistance is the between distance of two paragraphs in Pixels in the yAxis
            ParToParDistance = space*yStep + (space+1)*(yStep*textBoxInfo["lineSpacing"])
            space = 0 #re-initialize space between paragraphs
            ParToParDistanceList.append(ParToParDistance)
    ParToParDistanceList.append(0)#This value is needed for the last paragraph in order to be in bounds!!!!!!!
    ####------------------------------------------------------------------------------------------
    ####------------------------------------------------------------------------------------------
    ####------------------------------------------------------------------------------------------
    ####------------------------------------------------------------------------------------------
                
    for par in paragraphsList:
        parLinesWidth = []
        if not par: #If paragraph line is empty and has no text
            space +=1
            print("Paragraph line is empty\n")
            #We need the height of the font in get it with space?
        else:
            print("Paragraph line is full")
            par_counter+=1
            paragraphLines = 1 #initialize paragraph lines counter
            realParTextLength = get_pil_text_size(par,textBoxInfo["fontSize"], textBoxInfo["fontName"])[0]
            pptxParlength = realParTextLength + textBoxInfo["leftMargin"] + textBoxInfo["rightMargin"] # This is the real length because we have to consider the margins of the textBox
            if pptxParlength > shapeLoc["width"]:
                if (pptxParlength/shapeLoc["width"]) > int(pptxParlength/shapeLoc["width"]):
                    paragraphLines = int(pptxParlength/shapeLoc["width"])+1
            i=1 #i is like a line counter inside the paragraph
            #parLinesWidth = []
            while i <= paragraphLines:
                if i == paragraphLines:
                    print("realParTextLength",realParTextLength)
                    print("realParWidth * i",realParWidth * i)
                    print("realParWidth",realParWidth)
                    print("i",i)
                    parLinesWidth.append(realParTextLength - (realParWidth * (i-1)))
                else:
                    parLinesWidth.append(realParWidth)
                i+=1
            print("PARA:",parLinesWidth)

                        
            if par_counter==1 and space > 1: #This will handle the use case where before the firs par there are empty lines- empty paragraphs
                #MARGIN IN PIXELS margin =  yStep*textBoxInfo["lineSpacing"]
                parHeight = paragraphLines*yStep + (paragraphLines-1)*(yStep*textBoxInfo["lineSpacing"])
                #yDistance is the between distance of two paragraphs in Pixels in the yAxis
                ParToParDistance = space*yStep + (space+2)*(yStep*textBoxInfo["lineSpacing"]) #THIS IS space+2 because we have deleted the initial space from realYaxisStart
                paragraphHeightStart = paragraphHeightStart - ParToParDistance # stores first paragraph beginning
                nextParYstart = parHeight + ParToParDistanceList[par_counter]
                print("malakia")
                print("paragraphHeightStart:", paragraphHeightStart)
                print("ParToParDistance:", ParToParDistance)
                print("parHeight:", parHeight)
                print("yStep:", yStep)
                print("marginSpace:", (space+1)*(yStep*textBoxInfo["lineSpacing"]))
                space = 0 #re-initialize space between paragraphs
            else:
                #MARGIN IN PIXELS margin =  yStep*textBoxInfo["lineSpacing"]
                parHeight = paragraphLines*yStep + (paragraphLines-1)*(yStep*textBoxInfo["lineSpacing"])
                #yDistance is the between distance of two paragraphs in Pixels in the yAxis
                ParToParDistance = ParToParDistanceList[par_counter]# this is the par to par distance that was calculated before!
                paragraphHeightStart = paragraphHeightStart - nextParYstart # stores each paragraph height start
                print("paragraphHeightStart:", paragraphHeightStart)
                print("nextParYstart:", nextParYstart)
                print("ParToParDistance:", ParToParDistance)
                print("parHeight:", parHeight)
                print("yStep:", yStep)
                print("marginSpace:", (space+1)*(yStep*textBoxInfo["lineSpacing"]))
                nextParYstart = parHeight + ParToParDistance #calculates next paragraphs height start on y axis in pixels
                print("nextParYstartForNextLINE:", nextParYstart)
                space = 0 #re-initialize space between paragraphs
                        

            #change to correct resolution
            realXaxisStart = ScaleToScreenResolution(slideRes, realXaxisStart)
            realYaxisStart = ScaleToScreenResolution(slideRes, paragraphHeightStart)
            realParWidthPixel = ScaleToScreenResolution(slideRes, realParWidth)
            realHeight = ScaleToScreenResolution(slideRes, parHeight)
            yPixelStep = ScaleToScreenResolution(slideRes, yStep)#yPixelStep is not yStep because it will call it again and again and multiply it by scale var and deform the data that use ystep above
            Pixelmargin = ScaleToScreenResolution(slideRes, margin)#same as above for the margin!
            parLinesWidth = ScaleToScreenResolution(slideRes,parLinesWidth)

                        
                        
            #Get Occurrences- Duration of Focus
            print(parLinesWidth)
            OCnum = getGazeData(df, parLinesWidth, realXaxisStart, realYaxisStart, yPixelStep, Pixelmargin)
            print("RESULT OF getGazeData:",OCnum)

            shapeInfo = TextBox(realXaxisStart, realYaxisStart, realParWidthPixel, realHeight, "TEXT_BOX", 0, 1, OCnum, slideNum, par, parLinesWidth, Pixelmargin, yPixelStep, textBoxInfo["fontName"], textBoxInfo["fontSize"], textBoxInfo["leftMargin"], textBoxInfo["rightMargin"], textBoxInfo["topMargin"], textBoxInfo["botMargin"])#  Create object of class TextBox
            #shape_dict['Shape_' + str(shape_counter) + '_slide_' + str(slideNumber) + '_par_' + str(par_counter)] = shapeInfo# Insert Object to dictionary with key Slide_'#slide_number'
            #SOS SOS SOS SOS DO I NEED TO SAVE The shape to dict with dynamic names because every with only one shapeInfo we will only save there the last one on memory
            #ALSO SOS SOS SOS SOS WRITE TO REPORT THAT IT VIEWS EAXH PARAGRAPH AS A SEPARATE TEXTBOX ELEMENT EVEN IF IT IS IN ONE TEXTBOX.
            shapeInfo.getValues()

            #To Do:
            #SOS prepei na exoume kai use case stin periptwsi pou bazei to text grouped eite sthn mesh eite sta dexia anti gia to default aristera  
            #SOS OTAN EXOUME TWRA ESTW TITLO ME TOPOTHETISI STHN MESH OPWS EINAI O KWDIKAW TWRA THA KANEI STA MAP TO WIDTH APO THN ARISTERH PLEURA KAI OXI STHN MESH TO KEIMENO OPOTEW PROBLEM!!!     
            #Complete the function to get data from the dataframe using x and y and get the timestamps and occurences from the eyetracking.py
            #Open subproccess for the tobiistream.exe and create a interupt to stop it so it can continue later on  with a flag for the next slide.
            #Code Refinement and remove unessesary code from classes.py

def picture(slideRes, slideNum, picture_shape, df):

    image_loc = locateShape(picture_shape, slideRes)
    print("This is an image:",image_loc)

    width = [image_loc["width"]]
    OCnum = getGazeData(df, width, image_loc["xAxisStart"], image_loc["yAxisStart"], image_loc["height"], 0)
    print("RESULT OF getGazeData for Picture:",OCnum)

    img_blob = picture_shape.image.blob
    img_extensionType = picture_shape.image.ext

    imageInfo = Image(image_loc["xAxisStart"], image_loc["yAxisStart"], image_loc["width"], image_loc["height"], "PICTURE", 0, 1, OCnum, slideNum, img_blob, img_extensionType)
    imageInfo.getValues()
    
    #Check also https://stackoverflow.com/questions/58804348/convert-blob-image-to-png-jpg-with-python
    #CHECK https://stackoverflow.com/questions/62864082/copy-slide-with-images-python-pptx AND SOS LOOK LAST COMMENT ABOUT UNIQUINESS
    # image_loc = {"xAxisStart": xAxis, "yAxisStart": yAxis, "width": width, "height": height}
    # I could add image rotation here????????????????

def gatherData(data_frame, powerpointPath):

    prs = Presentation(powerpointPath)

    for slide in prs.slides:
        width = prs.slide_width/9525
        height = prs.slide_height/9525
        slideResolution = [width, height]
        slideNumber = prs.slides.index(slide)

        slideInfo = Slide(width, height, slideNumber)#  Create object of class Slide
        slide_dict['Slide_' + str(prs.slides.index(slide))] = slideInfo# Insert Object to dictionary with key Slide_'#slide_number'

        shape_counter = 0
        group_counter = 0
        for shape in slide.shapes:
            print("Type of shape", shape.shape_type)#   Type of shape!
            shape_counter+=1    #groups are seen as one shape
            #For Textboxes
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:    
                textbox(slideResolution, slideNumber, shape, data_frame, None) #None because of no associated group
                    

            #print("TEXTBOXES INFO\n")
            #print(shape_dict)
            
            #IF SHAPE IS A PICTURE
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture(slideResolution, slideNumber, shape, data_frame, None) #None because of no associated group
                print("This is an image")

            #IF SHAPE IS A GROUP OF OTHER SHAPES(TEXT_BOX AND PICTURE)
            #!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!
            #!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!
            #!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!

            #SOS WE NEED TO CONNECT THE PICTURES WITH THE TEXT WITH SOME WAY MAYBE A SPECIFIC UNIQUE KEY
            #I COULD CREATE A STORED PROCEDURE OR TRIGGER IN MYSQL DATABASE THAT WILL HANDLE THAT DATA.
            #MAYBE DO THIS
            #CREATE A NEW COLUMN IN THE TABLE OF THE SHAPES CALLED "CORELLATION_ID"
            #AUTO TO ID THA DEIXNEI PALI PISW STO ID TOU OBJECT(EIKONA H KEIMENO POY EISIXTH PRWTO)
            #OPOTE SE KA8E MIA EISAGWGH SUB-SHAPE TO EPOMENO PRAGMA EINAI CALL STHN DATABASE GIA TO ID TOU WSTE AUTO NA PAEI SAN CORELLATION ID STA EPOMENA

            #!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!
            #!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!
            #!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                print("This is a group")
                groupName = group_counter
                group_counter+=1
                #this will work for simple 1st level grouping it wont iterate in the more inside groups!!!
                #if we want to iterate more check the bellow link!!!!
                for sub_shape in shape.shapes:
                    if sub_shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX: 
                        print("This is a TextBox")
                        textbox(slideResolution, slideNumber, shape, data_frame, groupName)
                    if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        print("This is an Image")
                        picture(slideResolution, slideNumber, shape, data_frame, groupName)
            # https://stackoverflow.com/questions/51701626/how-to-extract-text-from-a-text-shape-within-a-group-shape-in-powerpoint-using

#prs.save('test2.pptx')

#test dictionary with slide classes.
print(slide_dict)

test = slide_dict.values()
for class_elem in test:
    class_elem.getSlideNum()



#TO DOOOOOOOOOOOOO

#check ParToParDistance it saves wrong the data maybe use a list? or reverse the current par list!!!!

#ABOUT TIMESTAMP MAYBE A FUNCTION TO SEE WHEN THE USER EYES GET OUT OF THE X,Y COORDINATES OF THE OBJECT AND THEN USE THIS TIMESTAMP AS MAX,
#THEN CHECK WHEN AND IF HE REENTERS THIS OBJECT AND CRAETE A NEW MIN-MAX AFTER HE SEES OUTSIDE SO WHEN HE STOPS THE DIFF MAX-MIN WILL BE ADDED TO THE FINAL TIME TO SEE
#HOW MUCH TIME THE USER HAS FOCUSED ON AN OBJECT