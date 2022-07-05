import os
from re import S
import MySQLdb
from pptx import Presentation
from pptx.util import Cm, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import ImageFont
from PIL import Image as pil_img
import tkinter
import MySQLdb
import mysql.connector
##Create Synopsis

def folderCheck(name, presentationName): #the data will be a list or a dict containing the objects that we need to insert to a powerpoint
    #CREATE SAVE FOLDER!
    path=os.getcwd()+"\\synopsis\\"+name#we could use the session["name"] which is the userID for anonymity but because of the hash string it might be unreachable by the os beacuse of it's characters so we will use the username!
    #if file not exists create it!
    if not os.path.exists(path):
        os.mkdir(path)
        print("Directory " , path ,  " Created ")
        print("User synopsis folder was Created!")
        return True
    else:
        print("Directory " , path ,  " already exists")
        print("User synopsis folder already exists!")
        filePath = path + "\\" + presentationName + ".pptx"
        #file_exists = os.path.exists(filePath)
        if not os.path.exists(filePath):
            print("No old pptx synopsis exist proceed normal.")
            return True
        else:
            print("Attetion old pptx synopsis found.")
            print("Deleting old to be replaced by new one.")
            os.remove(filePath)
            if not os.path.exists(filePath):
                print("File deleted successfully")
                return True
            else:
                print("File deletion FAILED!")
                return False

def select_prs_data(prsDataList, percentageValue):
    selectionVar = percentageValue/100
    data_size= len(prsDataList)
    remove_elem = selectionVar*data_size    #IF IT IS <1 THE DO IT 1
    if remove_elem < 1:
        remove_elem = 1 #because of very small sample we will at least show one
    #keep_data = int(data_size-remove_elem) # Maybe we need +1
    #SOS SOS SOS SOS SOS SOS SOS MAYBE WE ONLY NEED THE remove_elem

    print("TESTINGGGGGGG PERCENTACE GET")
    print("OLD DATA LIST:",prsDataList)
    #print("NEW DATA LIST:",prsDataList[0:keep_data])
    print("Percentage:",selectionVar)
    print("data_size:",data_size)
    print("remove:",remove_elem)

    return prsDataList[0:int(remove_elem)]  #correct one maybe we could also do a round up if we want!
    
    #SOS SOS FOR TESTING ONLY RETURN ALL NO REMOVAL
    #return prsDataList

def splitData(dataList):
    textbox_list = []
    picture_list = []
    positionCounter=0
    for dict in dataList:
        if(dict['objectCategory']=="TEXT_BOX"):
            dict['positionCounter'] = positionCounter #Add a new key to the dictionaries that will say their correct position based on previous sorting that happened on the DataBase
            textbox_list.append(dict)
            positionCounter = positionCounter + 1

        elif(dict['objectCategory']=="PICTURE"):
            dict['positionCounter'] = positionCounter #Add a new key to the dictionaries that will say their correct position based on previous sorting that happened on the DataBase
            picture_list.append(dict)
            positionCounter = positionCounter + 1

    return textbox_list, picture_list

def slideCalculation(dataList):
    #   Based on simple 1280*720P default slide of powerpoint
    #   If we want we could add more functionality by searching the slides table of the DB and getting custom slides based on the elements
    #   of the slide that we have on synopsis creation!
    #   for our experiments we used only default(1280*720p) slides in order not to overcomplicate things so the pptx images and text will fit on 1280p slides
    #   In other resolutions we should check for each element especialy pictures if it will fit in the default slide or not!

    #root = tkinter.Tk()
    #screen_res = [root.winfo_screenwidth(), root.winfo_screenheight()]
    #scaleMultiplier = screen_res[0]/1280    #this should be changed base on the slide resolution 1.5
    #we have to divide the pixels by the scale multiplier to see where the elements are not when the screen up-scales but where they are inside the slide preview not while presenting
    #availabelSlideVolume = 1280*720
    availabelSlideVolume = 720
    neededSlideVolume = 0
    counterChange = 1
    maxElementPerSlide = 0
    slideChangeList=[] # slideChangeList starts element counting from 1 not 0!
    while_flag=0 #while flag
    overflowFlag = False
    first_flag = True
    lastListElement=0   #counts the elements of the list so that we can append the last remaining elements to slideChangeList
    listLength = len(dataList)
    for dict in dataList:
        scaleMultiplier = dict['scaleMultiplier']
        print("\nscaleMultiplier!!!scaleMultiplier=",scaleMultiplier)

        lastListElement+=1
        while_flag+=1
        maxElementPerSlide+=1

        if(dict['objectCategory']=='TEXT_BOX'):
            elementWidth = dict['textBoxWidth_no_margin']/scaleMultiplier
            elementHeight = dict['textBoxHeight_no_margin']/scaleMultiplier + 42   #5 pixles of the 30 are because of the top margin of the textboxes which is 0.13cm # +10 for safety restriction
            print("TEXTBOX EL")

        elif(dict['objectCategory']=='PICTURE'):
            elementWidth = dict['pictureWidth']/scaleMultiplier
            elementHeight = dict['pictureHeight']/scaleMultiplier + 42   #5 pixles of the 30 are because of the top margin of the textboxes which is 0.13cm # +10 for safety restriction
            print("PICTURE EL")

        print("For loop")
        print("lastListElement ",lastListElement)
        print("while_flag ",while_flag)
        print("maxElementPerSlide ",maxElementPerSlide)
        #use a counter inside the for loop and compare when the needed_slide_volume goes above one slide(availabelSlideVolume) to know when to choose new slide
        #so how many elements we will have on one page
        #also here rises a problem with one more check not to split the grouped elements!!!

        #curElementVolume = elementWidth*elementHeight
        curElementVolume = elementHeight
        neededSlideVolume = curElementVolume + neededSlideVolume

        print("\n ")
        print("neededSlideVolume ",neededSlideVolume)
        print("availabelSlideVolume*counterChange ",availabelSlideVolume*counterChange)
        #Check for overflow on a slide
        if (neededSlideVolume > availabelSlideVolume*counterChange):
            overflowFlag = True
            counterChange+=1
            print("OVERFLOW")
            print("Element that overflowed")
            print(dict)
            print("\n")
            #Check the overflowing element and handle it accordingly if it is in a group or if it is not!
            if dict['groupID'] is None:

                #if it is the first element that overflowed insert it in the first slide
                if(first_flag and (curElementVolume > availabelSlideVolume)): #(curElementVolume > availabelSlideVolume is added to the if as a flag because if we can put more elements in the first and then it overflows it should go to the else
                    first_flag = False
                    slideChangeList.append(maxElementPerSlide)
                    print("Overflow")
                else:
                    slideChangeList.append(maxElementPerSlide-1) # -1 because the maxElementPerSlide will be put to the next slide as it won't fit to the previous!
                print("check slideChangeList1",slideChangeList)
                maxElementPerSlide = 1 #Re-initialize it for new slide element counter, it is one because of the overflowed element
                
                print("lastListElement:",lastListElement)
                print("listLength:",listLength)
                #if lastListElement == listLength:#If it is the Last Element of the list insert in it. but we have to take a case if it is not the last element to save the maxElementPerSlide
                #    slideChangeList.append(maxElementPerSlide)
                #    #set it False because if overflow wont happen again we need to save the maxElementPerSlide of this slide
            else:
            # check if previous pptx dict element group id matches this one so it would be too added to the next slide together as one
            # also if other groups exist before this we have to change the counterChange, neededSlideVolume and maxElementPerSlide!
                i=0
                groupFlag = False
                while i <= while_flag-1:
                    tempDict = dataList[i]

                    #If there are group elements of the same group as this overflowed element before it put em together to a new slide! 
                    if(tempDict['slide_number'] == dict['slide_number'] and tempDict['groupID'] == dict ['groupID']):
                    # Because grouped Elements are put sequential it meants that if dataList[i] is in the same slide number and group as the one that will be put on the next page
                    # that all the elements from i to maxElementPerSlide will be on the same slide as they are grouped!
                        groupFlag = True
                        groupList = dataList[i:while_flag]#make a list of the sequential matched groups!
                        print("GROUP ELEMETS LIST",groupList)
                        groupVolume = 0

                        for groupElement in groupList:
                            
                            if(groupElement['objectCategory']=='TEXT_BOX'):
                                print("Group ELEMENT",groupElement)
                                groupElementWidth = groupElement['textBoxWidth_no_margin']/scaleMultiplier
                                groupElementHeight = groupElement['textBoxHeight_no_margin']/scaleMultiplier + 30   #5 pixles of the 30 are because of the top margin of the textboxes which is 0.13cm # +10 for safety restriction
                                print("TEXTBOX GROUP ELEMENT")

                            elif(groupElement['objectCategory']=='PICTURE'):
                                print("Group ELEMENT",groupElement)
                                groupElementWidth = groupElement['pictureWidth']/scaleMultiplier
                                groupElementHeight = groupElement['pictureHeight']/scaleMultiplier + 30   #5 pixles of the 30 are because of the top margin of the textboxes which is 0.13cm # +10 for safety restriction
                                print("PICTURE GROUP ELEMENT")
                            groupVolume = groupVolume + groupElementHeight
                            print("groupVolume",groupVolume)

                        #Check if group size exceeds slide size
                        if(groupVolume > availabelSlideVolume):
                            print("******ERROR*********\nGROUP SIZE IS BIGGER THAN THE SLIDE SIZE!!!")
                            i = while_flag #to break the while
                        #else:
                        #here should be all the rest! 134-145

                        #add  these elements to neededSlideVolume for next slide calc
                        neededSlideVolume = neededSlideVolume + groupVolume #it's "+" instead of "-" because we get to a new slide and we don't care about the old one!
                        print("old maxElementPerSlide=",maxElementPerSlide)
                        print("len(groupList)",len(groupList))
                        i = while_flag #to break the while   
                        print("FIX ME SOS")
                        print("maxElementPerSlide",maxElementPerSlide)
                        print("len(groupList)",len(groupList))
                        maxElementPerSlide_fix = maxElementPerSlide - len(groupList) #because the maxElementPerSlide has also the +1 from the initialization which is from the group list
                        #SOS IF THE OVERFLOW OF GROUPS ARE TWO CONTINUED OF THE SAME GROUP THEN  maxElementPerSlide_fix=0=>ERROR CAUSE IT CANT FIT IN SLIDE SIZE
                        print("New correct maxElementPerSlide_fix=",maxElementPerSlide_fix)

                        #insert to slideChangeList the elements the will be on the slide
                        slideChangeList.append(maxElementPerSlide_fix)#remove the groups that have to go with the overflowing element
                        print("check slideChangeList2",slideChangeList)
                        maxElementPerSlide = maxElementPerSlide-maxElementPerSlide_fix #to continue correctly for the next slide also we need -1
                        
                        #if lastListElement == listLength:#If it is the Last Element of the list insert in it.
                        #    slideChangeList.append(maxElementPerSlide)
                    
                    i+=1 # increase counter i for while loop escape
                print("WHILE_END")
                #If there is not a group matching before this overflowed element simply change slide and put it there
                if not groupFlag:
                    print("No group elements found before the overflow element that are associated with it")
                    slideChangeList.append(maxElementPerSlide-1)
                    print("check slideChangeList3",slideChangeList)
                    maxElementPerSlide = 1 #Re-initialize it for new slide element counter
                    
                    #if lastListElement == listLength:#If it is the Last Element of the list insert in it.
                    #    slideChangeList.append(maxElementPerSlide)

            #ELSE GROUP OVERFLOW ID END

        #IF OVERFLOW END     
       
    #for end
    if not overflowFlag or lastListElement == listLength:   #If overflowFlag is false everyElement we want to insert fits on one page so the slideChangeList would be empty\ or the last sum of elements fits ok in the last page with no overflow
        slideChangeList.append(maxElementPerSlide)
        print("check slideChangeList4",slideChangeList)


    #if counterChange-1 == 0: #because it should be initialized as 0 but we initializes it as 1 to be able to increase available height
    #    counterChange = 1

    #slideChangeList  is a list with integers to show when to change elements on new slide, counterChange equals to how many slides we will need for the synopsis!
    return slideChangeList, counterChange
    #MAYBE counterChange IS OBSOLETE AND WE BETTER USE len(slideChangeList)

#This function will find the needed groups of a pptx element if it has any and put them sequentially together to be later used at slideCalculation function and to the final powerpoint synopsis creation
def groupCheck(dataList, pptxName):
    duplicateListCheck=[]
    groupCorrectedList=[]
    #Only get grouped elements!
    #The select queries are ORDER BY slide_number ASC, object_counter ASC to insert them as correctly as possibly
    #TextboxElements
    conn_object=mysql.connector.connect(host="localhost", user="root", passwd="", db="ez-synopsis")
    cursor = conn_object.cursor(dictionary=True)#(MySQLdb.cursors.DictCursor)
    cursor.execute('SELECT objectCategory, slide_number, textBoxWidth_no_margin, textBoxHeight_no_margin, parText, eachLineWidth, lineSpacer, lineSizeY, fontName, fontSize, scaleMultiplier, groupID FROM textbox_elements WHERE pptx_name = %s AND groupID IS NOT NULL ORDER BY slide_number ASC, object_counter ASC', (pptxName,))
    textboxGroups= cursor.fetchall()
    cursor.close()
    #PictureElements
    #conn_object=mysql.connector.connect(host=hostname, user=username, passwd=password, db=database)
    cursor = conn_object.cursor(dictionary=True)
    cursor.execute('SELECT objectCategory, slide_number, pictureWidth, pictureHeight, imageLoc, scaleMultiplier, groupID FROM picture_elements WHERE pptx_name = %s AND groupID IS NOT NULL ORDER BY slide_number ASC, object_counter ASC', (pptxName,))
    pictureGroups= cursor.fetchall()
    cursor.close()
    #get all groups associated with this list of dicts
    #delete later used groups
    #put all groups sequential
    #fix the positionCounter  to be based on the new list of dict
    #we need two database calls one for "TEXT_BOX" AND ONE FOR "PICTURE"
    for pptxElement_dict in dataList:
        print("GROUP ID IS:",pptxElement_dict['groupID'])
        if pptxElement_dict['groupID'] is None: # OR IF IT IS "NULL" as an string?
            print("slide_number:",pptxElement_dict['slide_number'])
            print("None as GROUP ID")
            groupCorrectedList.append(pptxElement_dict)
        else:
            #Check for duplicate!
            if [pptxElement_dict['slide_number'],pptxElement_dict['groupID']] in duplicateListCheck:
                print("This is a Duplicate of an already implemented group so it is not needed again!")
            else:
                #Check TextBoxGroups and PictureGroups For additional group info also register what was the groupID and the slideNumber of the element
                for textboxGroupElem_dict in textboxGroups:
                    if pptxElement_dict['slide_number'] == textboxGroupElem_dict['slide_number'] and pptxElement_dict['groupID'] == textboxGroupElem_dict['groupID']:
                        textboxGroupElem_dict['positionCounter'] = pptxElement_dict['positionCounter']
                        groupCorrectedList.append(textboxGroupElem_dict)

                for pictureGroupElem_dict in pictureGroups:
                    if pptxElement_dict['slide_number'] == pictureGroupElem_dict['slide_number'] and pptxElement_dict['groupID'] == pictureGroupElem_dict['groupID']:
                        pictureGroupElem_dict['positionCounter'] = pptxElement_dict['positionCounter']
                        groupCorrectedList.append(pictureGroupElem_dict)

                #duplicate check list creation!
                duplicateListCheck.append([pptxElement_dict['slide_number'],pptxElement_dict['groupID']])

    #Maybe later add a sorting function based on the object_counter for these elements! it's not vital now but it could be a good feature other than just using the order by at the select queries!
    #ALSO HAVE TO TAKE CHANCE ON GROUP OVERLAPPING SO IT WONT SHOW TWICE THE SAME THING! 
    return groupCorrectedList

def synopsis_creation(username, presentationName, synopsisList): #the data will be a list or a dict containing the objects that we need to insert to a powerpoint
    print("username:",username)
    print("presentationName:",presentationName)
    checkFolder = folderCheck(username, presentationName)
    finalDictList = groupCheck(synopsisList, presentationName)
    synopsisHandler = slideCalculation(finalDictList) #slideChangeList, counterChange
    overflowCounter = synopsisHandler[1] 
    creationPath=os.getcwd()+"\\synopsis\\"+username+"\\"+presentationName + "_" + username + ".pptx"
    #######SOS SOS SOS SOS SOS SOS SOS SOS SOS TO DO!!!!
    #GIA NA MHN KANW POLLLA SELECT KAI APO TA DUO  TABLES GIA TEXTBOX KAI PICTURES
    #NA BALW STO SYNOPSIS_DATA TABLE ENA NEO COLUMN POU NA LEEI THN PROELEUSH TOU- EIDOS TOU TEXT  H PICTURE
    # KAI ANALOGA NA KALW TA SELECT GIA TA DEDOMENA GIA THN DHMIOURGIA NEOU PPTX!
    #objectCategory
    
    list_idx=0
    slice_idx = synopsisHandler[0][0] #initial slice.
    synopsisTemplateLoc = os.getcwd() + "\\templates\\Powerpoint_Synopsis_Templates\\" + "template.pptx"
    prs = Presentation(synopsisTemplateLoc)
    blank_slide_layout = prs.slide_layouts[6] #no6 is blank slide
    #1 pixel= 0.0264583333 inches
    pixelToInch = 0.01041666666
    #SOS SOS SOS SOS SOS SOS SOS -----------------------------------------------------------***************************************************//////////////////////////////
    #ADD SCALE MULTIPLIER TO TEXTBOX_ELEMENTS IN DATABASE
    #BECAUSE WE NEED IT FOR CORRECT INSERTION FOR FINAL SYNOPSIS
    #FOR NOW I WILL ADD IT HARD CODED FOR RESOLUTION 1920*1080 AND SLIDE RES 1280*720p
    #SO THE SCALE MULTIPLIER WOULD BE 0.5 in powerpoint this is 1.5lines.
    #SOS SOS SOS SOS SOS SOS SOS -----------------------------------------------------------***************************************************//////////////////////////////
    textframe_flag = False
    #TextFrameYstart = 0#could be 0
    if(checkFolder):#Proceed with creating the synopsis file
        #if overflowCounter == 0: #if it is 0 it menas all elements fin in one slide!
        #    overflowCounter = 1

        for slide_num in range(overflowCounter):
            print("SOS SOS synopsisHandler[1]:",overflowCounter)
            #add slide and basic info
            TextFrameYstart = 0
            slide = prs.slides.add_slide(blank_slide_layout)
            left = Inches(0.20)
            top = Inches(0)
            width = Inches(13) # 13 inches so that the textbox will be as big as the 1280p of the slide.
            height = Inches(0) #Inches(2)
            
            #textBox = slide.shapes.add_textbox(left, top, width, height)
            #tf = textBox.text_frame

            print("NEW SLIDE CREATED, PROCEED WITH INSERTING")
            #for pptxElementCounter in range(len(synopsisHandler[0])): # + 1 exist here because we need to take into account the first slide because synopsisHandler[0] is appened only on overflow
            #    #MAYBE INSTEAD OF  len(synopsisHandler[0]) + 1) WE COULD USE len(synopsisHandler[1]))
            #
            #    print("LENGTH OF synopsisHandler[0]:",len(synopsisHandler[0]))
            elementListSlice = finalDictList[list_idx:slice_idx]
            #elementListSlice.reverse()#reverse the list. The reverse() method doesn't return any value. It updates the existing list.

            for pptxElement_dict in elementListSlice: #list slicing! Slice finalDictList based on slideChange_list, reverse used because of the way the pars are inserted so that fistly the most important text will be in first line!
                print("Insert into this slide")
                print("\n")
                print("synopsisHandler[0]",synopsisHandler[0])
                print("synopsisHandler[1]",synopsisHandler[1])
                print("list_idx: ",list_idx)
                print("slice_idx: ",slice_idx)

                print("\n")

                #if it is text
                if pptxElement_dict['objectCategory']=="TEXT_BOX":
                    print("Inserting TextBox")

                    #If it has inserted a picture before we need to add a new textFrame
                    if (textframe_flag):
                        #create new textBox
                        
                        top = Inches(TextFrameYstart*pixelToInch)
                        #new top  is the end of the last par of the old textframe
                        height = Inches((pptxElement_dict['textBoxHeight_no_margin']+10)*pixelToInch) #+10 is to add again a margin of top and bot of 0.13cm
                        textBox = slide.shapes.add_textbox(left, top, width, height)
                        tf = textBox.text_frame
                        print("New TextFrame created!")
                        tf.text = pptxElement_dict["parText"]
                        tf.paragraphs[0].font.name = 'Calibri'
                        #tf.paragraphs[0].font.size.pt=12
                        tf.fit_text()
                        #par = tf.add_paragraph()
                        #par.text = pptxElement_dict["parText"]

                        print("Par Inserting:",pptxElement_dict["parText"])
                        #SOS SOS SOS SOS SOS SOS SOS SOS SOS SOS SOS SOS SOS SOS
                        eachLineWidth = pptxElement_dict['eachLineWidth'].split(",")
                        if isinstance(eachLineWidth, list): #If it is a LIST
                            print("Is Instance")
                            print("TextFrameYstart=",TextFrameYstart)
                            print("eachLineWidth",len(eachLineWidth))
                            print("pptxElement_dict['lineSizeY']",pptxElement_dict['lineSizeY']/pptxElement_dict['scaleMultiplier'])
                            #textYheight = (pptxElement_dict['lineSizeY']/pptxElement_dict['scaleMultiplier'])*len(eachLineWidth) + pptxElement_dict['lineSpacer']*len(eachLineWidth)
                            textYheight = pptxElement_dict['lineSizeY']*len(eachLineWidth) + pptxElement_dict['lineSpacer']*len(eachLineWidth)
                            TextFrameYstart = TextFrameYstart + textYheight + 15   #5 pixles of the 30 are because of the top margin of the textboxes which is 0.13cm
                            print("TextFrameYstart AFTER=",TextFrameYstart)
                        else:
                            print("TextFrameYstart=",TextFrameYstart)
                            print("pptxElement_dict['lineSizeY']",pptxElement_dict['lineSizeY'])
                            #textYheight = pptxElement_dict['lineSizeY']/pptxElement_dict['scaleMultiplier'] + pptxElement_dict['lineSpacer']
                            textYheight = pptxElement_dict['lineSizeY'] + pptxElement_dict['lineSpacer']
                            TextFrameYstart = TextFrameYstart + textYheight + 15   #5 pixles of the 30 are because of the top margin of the textboxes which is 0.13cm
                            print("TextFrameYstart AFTER=",TextFrameYstart)
                        #If the text is too long and is cutted by the textFrame then to go bellow then it is correct else it's false
                        # WE could use instead of pptxElement_dict['textBoxHeight_no_margin'] pptxElement_dict['lineSizeY']* len(', '.join(list(map(str, class_dict['EachLineWidth']))))
                        #Also think if we could use the original textBox! AS WE HAVE THE VALUES FROM THE DATABASE SO INSTEAD OF ADDING PARAGRAPHS WE COULD ADD TEXTFRAMES WITH SPECIFIC VALUES AS INTENDED
                        #SOS SOS SOS SOS SOS SOS SOS SOS SOS SOS SOS SOS SOS SOS

                        textframe_flag = False #Re Initialize Flag

                    #No picture was added right before the new text we need so continue inserting to par
                    else:
                        print("ELSE Par Inserting:",pptxElement_dict["parText"])

                        top = Inches(TextFrameYstart*pixelToInch)
                        height = Inches((pptxElement_dict['textBoxHeight_no_margin']+10)*pixelToInch) #+10 is to add again a margin of top and bot of 0.13cm
                        textBox = slide.shapes.add_textbox(left, top, width, height)
                        tf = textBox.text_frame
                        tf.text = pptxElement_dict["parText"]
                        tf.paragraphs[0].font.name = 'Calibri'
                        #tf.paragraphs[0].font.size.pt=12
                        tf.fit_text()
                        #par = tf.add_paragraph()
                        #par.text = pptxElement_dict["parText"]
                        eachLineWidth = pptxElement_dict['eachLineWidth'].split(",")
                        if isinstance(eachLineWidth, list): #If it is a LIST
                            print("Is Instance")
                            print("TextFrameYstart=",TextFrameYstart)
                            print("eachLineWidth",len(eachLineWidth))
                            print("pptxElement_dict['lineSizeY']",pptxElement_dict['lineSizeY']/pptxElement_dict['scaleMultiplier'])
                            #textYheight = (pptxElement_dict['lineSizeY']/pptxElement_dict['scaleMultiplier'])*len(eachLineWidth) + pptxElement_dict['lineSpacer']*len(eachLineWidth)
                            textYheight = pptxElement_dict['lineSizeY']*len(eachLineWidth) + pptxElement_dict['lineSpacer']*len(eachLineWidth)
                            TextFrameYstart = TextFrameYstart + textYheight + 15   #5 pixles of the 30 are because of the top margin of the textboxes which is 0.13cm
                            print("TextFrameYstart AFTER=",TextFrameYstart)
                        else:
                            print("TextFrameYstart=",TextFrameYstart)
                            print("pptxElement_dict['lineSizeY']",pptxElement_dict['lineSizeY'])
                            #textYheight = pptxElement_dict['lineSizeY']/pptxElement_dict['scaleMultiplier'] + pptxElement_dict['lineSpacer']
                            textYheight = pptxElement_dict['lineSizeY'] + pptxElement_dict['lineSpacer']
                            TextFrameYstart = TextFrameYstart + textYheight + 15   #5 pixles of the 30 are because of the top margin of the textboxes which is 0.13cm
                            print("TextFrameYstart AFTER=",TextFrameYstart)                            

                #if it is picture
                elif pptxElement_dict['objectCategory']=="PICTURE":
                    #AN EINAI EIKONA THA THELOUME ENA SWITCH-FLAG COUNTER (0-1) WSTE STHN SYNEXEIA AN DIABASEI TEXTBOX STO IDIO SLIDE NA BALEI WS EIDODO NEW TEXT FRAME
                    #WSTE NA MPORESEI NA BALEI EKEI MESA TA YPOLOIPA PARS GIA NA EXEI MIA LOGIKH SEIRA!
                    # TO NEO TEXT FRAME THA MPEI ME BASH TA YSTEP TWN PROHGOUMENWN PARS SAN SYNOLO.
                    print("Inserting Picture")
                    print("cur workinf dir:",os.getcwd()+pptxElement_dict['imageLoc'])

                    pic_width = (pptxElement_dict['pictureWidth']/pptxElement_dict['scaleMultiplier']) * pixelToInch    #scale remove and convert to inches measurement
                    pic_height = (pptxElement_dict['pictureHeight']/pptxElement_dict['scaleMultiplier']) * pixelToInch
                    top = Inches(TextFrameYstart*pixelToInch)

                    #Insert picture to slide
                    print("Pic Inserting:",pptxElement_dict['imageLoc'])
                    pic = slide.shapes.add_picture(os.getcwd()+pptxElement_dict['imageLoc'], Inches(1), top, Inches(pic_width), Inches(pic_height)) #left #top distance 
                    
                    TextFrameYstart = TextFrameYstart + pptxElement_dict['pictureHeight']/pptxElement_dict['scaleMultiplier'] + 25   #5 pixles of the 30 are because of the top margin of the textboxes which is 0.13cm
                    
                    textframe_flag = True   #This variable works as a flag to show that on the next pptx element insertion
                                                #The previous element that was added was an image, so we need a new textFrame to succeed continuity!


            #INITIALIZE THE NEXT SLICING
            list_idx = slice_idx 
            print("slide before if:",slide_num)
            if slide_num != overflowCounter - 1 and (len(synopsisHandler[0]) > 1) : #this if acts as an index out of bounds protection for the next line because of "slice+1"
                print("list_idx",list_idx)
                print("slide",slide_num)
                print("slice_idx old",slice_idx)
                print("LENGTH OF synopsisHandler[0]",len(synopsisHandler[0]))
                slice_idx = synopsisHandler[0][slide_num+1] + list_idx
                #slice_idx = synopsisHandler[0][slide_num+1] + list_idx# OVERFLOW ERROR
                print("sliceStart:",list_idx)
                print("slice_idx new",slice_idx)

                

        print("Elements Insertion Finished")

        prs.save(creationPath)# save the final synopsis file


        # https://python-pptx.readthedocs.io/en/latest/user/text.html
        # https://www.geeksforgeeks.org/creating-and-updating-powerpoint-presentations-in-python-using-python-pptx/
        # https://python-pptx.readthedocs.io/en/latest/user/slides.html
        # https://stackoverflow.com/questions/44275443/python-inserts-pictures-to-powerpoint-how-to-set-the-width-and-height-of-the-pi
        
        #https://pythonprogramming.altervista.org/inserting-an-image-in-powerpoint-with-python/
        #https://towardsdatascience.com/creating-presentations-with-python-3f5737824f61
        #https://stackoverflow.com/questions/62395983/how-to-create-a-text-shape-with-python-pptx
        #https://python-pptx.readthedocs.io/en/latest/dev/analysis/shp-autofit.html

        #SOS SOS SOS SOS MAYBE MOD THE slideCalculation based on the height of each element so to use a horizontal evaluation to be more precise!!!!
        #Fit text to textbox also increase text box width!

        #SOS FIX γραμματοσειρα στα textboxes να ειναι ίδια με αρψικα textboxes
        # SUPER SOS SOS STO TextFrameYstart an h grammatoseira einai variable(διαφορα μεγεθη) τα πολυ μικρά μεγέθη μπορει να κάνουν overlap γιατι γινονται resize se 18αρια αλλα το  Y υπολογίζεται με το παλιο linesizeY
        
        #SOS FIX ΝΑ ΒΑΛΩ ΚΑΙ ΤΟ LINESPACING ΜΕΣΑ ΣΤΟΝ ΥΠΟΛΟΓΙΣΜΟ ΤΟΥ  HEIGHT STO CREATE_SYNOPSIS
        #SOS INSERT SCALE MULTIPLIER TO DATABASE

    else:   #if check is False  we have to check for the fileName if it already exist then delete it and replace it with the new one
        print("Sheeesh! FAILURE ON FILE")
#http://127.0.0.1:5002/synopsis/bitconPresentation/running